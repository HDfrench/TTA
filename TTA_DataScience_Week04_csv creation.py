#-------------------------------------------------------------------------------
# Name:        TTA_Week_04_Pandas
# Purpose:     Working with Pandas library, manipulating data imported
#              from a json file and generating a pptx presentation
# Author:      Hacene Dramchini
#
# Created:     20/05/2021
# Copyright:   (c) Meanie_3 2021
# Licence:     Free
#-------------------------------------------------------------------------------

#!/usr/bin/env python
import pandas as pd
import numpy as np
import json
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import os
from matplotlib import pyplot as plt


pr1 = Presentation()
#limit the display to 5 columns
pd.set_option('display.max_columns', 5)
plt.figure()
plt.rcParams.update({'font.size': 8})




# this is a class made to handle slides creation
# the content is closer to VBA than python, but it si to be expected as
# this manipulates an MS Office application
class sld:
    def __init__(self, data, dataset):
        def para_setting(p,  text, color = RGBColor(0,112,192),level =0, fsize = Pt(12), fbold = False, palign ="left", fital = False):
            p.font.color.rgb = color
            if fbold == True:
                p.font.color.rgb = RGBColor(55,96,146)
            p.text = text
            p.level = level
            p.font.size = fsize
            p.font.bold = fbold
            p.font.italic = fital
            if palign == "left":
               p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER

        self.layout = pr1.slide_layouts[data["layout"]]
        self.slide = pr1.slides.add_slide(self.layout)
        self.background = self.slide.background
        self.fill = self.background.fill
        self.fill.solid()
        self.fill.fore_color.rgb= RGBColor(254,234,218)
        self.title = self.slide.shapes.title
        para_setting(self.title.text_frame.add_paragraph(), data["title"], fbold = True, fsize = Pt(25))
        if data["layout"] != 8:
            for header, details in data["placeholders"].items():
                self.slide.shapes[int(header)].text_frame.clear()
                for plh, cnt in details.items():
                    para_setting(self.slide.shapes[int(header)].text_frame.add_paragraph(),
                    cnt["text"], level =cnt["level"], fsize = Pt(int(cnt["size"])), fbold = cnt["bold"],
                     palign =cnt["align"], fital = cnt["ital"])
        else:
            self.image = self.slide.placeholders[1]
            self.image.insert_picture(data["placeholders"][1])
            self.slide.placeholders[2].text_frame.clear()
            para_setting(self.slide.placeholders[2].text_frame.add_paragraph(), data["placeholders"][2])



# This function deals with importing the json file, manipulating the data,
# saving the dataset as a vsc file and returning the dataset for use in the powerpoint presentation
def deal_with_csv():
    # import destination specific data
    with open('Destinations.json') as json_dst:
        destinations = json.load(json_dst)
    # create a data set for destination specific data for reference and calculation
    frm_dest = pd.DataFrame(destinations)
    # create a copy of the destination specific data set that will be manipulated
    df_for_csv = frm_dest.copy(deep=True)
    #Cleaning the dataset to keep only columns that do not require calculations
    # (remove all useless columns)
    col_to_be_deleted = ["Pop. (M)",
                                "Risk (country)",
                                "Risk (city)",

                                "Rentals",
                                "All-inc (%)"]
    for item in col_to_be_deleted:
        del df_for_csv[item]

    # feedback level is going to be a new np array where:
        # ((     1 - (  number of accomodations / number of tourists )
        #                       +
        #       number of accommodations / population               )
        #                       *
        #               1 + (star average / 5)
        #                       *
        #               1 + (risk level / 5)
        #                       *
        #        1 + ( (risk level city - risk level country ) / 5)
        #                       *
        #             1 + (nb all-incl / nb accomodations)          )
        #                       *
        #              5

    # creation of a set of arrays to define ratios for calculation of feedback level
        # conversion of the all inclusive expressed as a percentage into a value
    nb_all_inclusive = np.array(frm_dest['Rentals'] /100
                                * frm_dest['All-inc (%)'])
        # ratio number of accommadations by number of tourists (<1)
    ac_per_tour_ratio = np.array([(1 - (frm_dest['Rentals'] /
                                    (frm_dest['Visits (M)']* 1000000)))])
        # ratio of accomodations compared to the population (<1)
    ac_per_pop_ratio = np.array([frm_dest['Rentals'] /
                                (frm_dest['Pop. (M)']* 1000000)])
        # ratio of all inclusive accommodation amongst all rentals (<1)
    allinc_ac_ratio = np.array([nb_all_inclusive/frm_dest['Rentals']])
        # conversion of the average accomodation star as a ratio ( < 1 )
    star_ratio = np.array([1 + (frm_dest['Stars (avg)']/5)])
        # conversion of risk ratio (1<= x <= 2)
    risk_ratio_city = np.array([1 +(frm_dest['Risk (city)']/5)])
        # city risk ratio compared to national risk level (-1 <= x <= 1)
    risk_ratio_national_comparison = np.array([( 1 + ((frm_dest['Risk (city)'] / frm_dest['Risk (country)'])/5))])

    # overall satisfaction levels array to be added to the
    satisfaction_array = np.round(((((ac_per_tour_ratio + ac_per_pop_ratio)/2)
                        * star_ratio * risk_ratio_city * risk_ratio_national_comparison
                        * allinc_ac_ratio) * 10), 2)

    # convert the array back to a list
    satisfaction_list = satisfaction_array.tolist()
    # add the list to the dataset df_for_csv in index position 2
    df_for_csv.insert(2, "Feedback", satisfaction_list[0], True)
    # convert all inclusive array converting percentage to number to a list
    all_inclusive_number = np.round(nb_all_inclusive.tolist(),0)
    #add the list to our dataset as integer
    df_for_csv.insert(5, "All-inc", pd.to_numeric(all_inclusive_number, downcast='integer'), True)

    # must include a way to calculate the rank of each city for each country (1, 2 or 3)
    # and add resulting column to dataframe at index 5

    # order by column country name and number of visits
    df_for_csv.sort_values(['Country', 'Visits (M)'], inplace=True)
    # create a ranking list
    rank_array = np.hstack([ np.array([3,2,1], dtype=int) for _ in range(15)]).tolist()
    # insert list in dataframe
    df_for_csv.insert(6, "Visit rank", rank_array, True)

    # remove the number of visit as not require anymore
    del df_for_csv['Visits (M)']

    # must include a way to check if the file is opened and ask user to close it
    df_for_csv.to_csv('csv_summary.csv', index=False)

    return frm_dest, df_for_csv

def main():
    #procedure to generate the 3 scatter plots
    def scatter_plot_image(dataset, colx, coly, name, display):
        myplot = dataset.plot(kind='scatter', figsize=(4.5, 3), title=display, x=colx, y=coly)
        fig = myplot.get_figure()
        fig.get_size_inches()
        fig.savefig(name, bbox_inches='tight', dpi=600)
    # generating the dataset
    ref_set, dataset = deal_with_csv()

    # generating  the graphics
    # creating a set of 3 scatter plots for analisys
    scatter_plot_image(dataset, 'All-inc', 'Feedback', 'Comp_feedback_allInc.jpg', 'Comparison Feedback / All-inclusive')
    scatter_plot_image(dataset, 'Feedback', 'Stars (avg)', 'Comp_feedback_stars.jpg', 'Average stars / All-inclusive')
    scatter_plot_image(dataset, 'All-inc', 'Stars (avg)', 'Comp_allInc_stars.jpg', 'Comparison All-inclusive / Average stars')

    # creating a bar chart
    # creating a copy of the dataset and sorting it by ascending feedback to see if ascending and stars are correlated
    temp_df = dataset[['City', 'Country', 'Feedback', 'Stars (avg)']].copy(deep=True).sort_values(by=['Feedback', 'Stars (avg)'])
#    temp_df.sort_values(by=['Feedback', 'Stars (avg)'])
    #changing the index to a dual index
    temp_df.set_index('City', 'Country', inplace =True)
    #generating the horizontal bar chart (too many values for a vertical one)
    plot_4 = temp_df.plot.bar()


    fig = plot_4.get_figure()
    fig.set_size_inches(6, 3.5)
    fig.savefig("Bar_chart.jpg", bbox_inches='tight', dpi=600)


    # setting the content of the powerpoint slides
    slide_content = {1: {"layout":0, "title": "Using Pandas",
            "placeholders": {1:
                {1: {"level": 0,
                    "align": "left",
                    "text": "TTA - HLA - Week 04 — Hacène Dramchini",
                    "size": 14,
                    "bold": False,
                    "ital": False}
                }
            }
        },
        2: {"layout": 3, "title": "Introduction",
            "placeholders":{1:
                {1: {"level": 0,
                    "align": "left",
                    "text": 'In order to create the data frame containing the desired data, we create a json file named "Destinations.json", holding the data about 3 cities in 15 countries.',
                    "size": 14,
                    "bold": False,
                    "ital": False},
                2: {"level": 0,
                    "align": "left",
                    "text": 'In order to create the data frame containing the desired data, we create a json file named "Destinations.json", holding the data about 3 cities in 15 countries.',
                    "size": 14,
                    "bold": False,
                    "ital": False},
                3: {"level": 0,
                    "align": "left",
                    "text": 'The statistics in that file come from various sources. Amongst others:',
                    "size": 14,
                    "bold": False,
                    "ital": False},
                4: {"level": 1,
                    "align": "left",
                    "text": 'Wikipedia',
                    "size": 12,
                    "bold": False,
                    "ital": False},
                5: {"level": 1,
                    "align": "left",
                    "text": 'hotelscombined.com',
                    "size": 12,
                    "bold": False,
                    "ital": False},
                6: {"level": 1,
                    "align": "left",
                    "text": 'Google',
                    "size": 12,
                    "bold": False,
                    "ital": False},
                7: {"level": 0,
                    "align": "left",
                    "text": 'And it looks like this (header shown):',
                    "size": 14,
                    "bold": False,
                    "ital": False},
                8: {"level": 0,
                    "align": "left",
                    "text": f'{ref_set.head()}',
                    "size": 10,
                    "bold": True,
                    "ital": False}
                } ,
            2:{1: {"level": 0,
                    "align": "left",
                    "text": 'This file was then manipulated using the pandas library to extract the relevant data, keeping some columns and deleting others and using the original data to get specific values.',
                    "size": 14,
                    "bold": False,
                    "ital": False},
                2: {"level": 0,
                    "align": "left",
                    "text": f'The final data set was saved as a csv file named \"csv_summary.csv\". This file has {dataset.shape[0]} rows and {dataset.shape[1]} columns.',
                    "size": 14,
                    "bold": False,
                    "ital": False},
                3: {"level": 0,
                    "align": "left",
                    "text": "Starts with",
                    "size": 14,
                    "bold": False,
                    "ital": False},
                4: {"level": 0,
                    "align": "center",
                    "text": f'{dataset.head(3)}',
                    "size": 10,
                    "bold": True,
                    "ital": False},
                5: {"level": 0,
                    "align": "left",
                    "text": "And ends with",
                    "size": 14,
                    "bold": False,
                    "ital": False},
                6: {"level": 0,
                    "align": "center",
                    "text": f'{dataset.tail(3)}',
                    "size": 10,
                    "bold": True,
                    "ital": False}
                }
            }
        },
        3: {"layout": 2, "title": 'Some insights 1/5',
            "placeholders":{1:
                {1: {"level": 0,
                    "align": "left",
                    "text": 'Here is a sample of the content of the resulting file:',
                    "size": 14,
                    "bold": False,
                    "ital": False},
                2: {"level": 0,
                    "align": "left",
                    "text": f'{dataset.iloc[2:8]}',
                    "size": 10,
                    "bold": True,
                    "ital": False},
                3: {"level": 0,
                    "align": "left",
                    "text": 'We can now look at some details:',
                    "size": 14,
                    "bold": False,
                    "ital": False},
                4: {"level": 0,
                    "align": "left",
                    "text": f"The mean number of all-inclusive accommodations within the selected destinations is {dataset['All-inc'].mean():.2f}",
                    "size": 14,
                    "bold": False,
                    "ital": False},
                5: {"level": 0,
                    "align": "left",
                    "text": 'Destination scores were calculated taking account several variables from the original data (see Slide 7 titles \"Some insights 5/5\")',
                    "size": 14,
                    "bold": False,
                    "ital": False},
                6: {"level": 0,
                    "align": "left",
                    "text": f"As a result, {tuple(dataset.loc[dataset[['Feedback']].idxmin(), 'City'])[0]} in {tuple(dataset.loc[dataset[['Feedback']].idxmin(), 'Country'])[0]} has the lowest rating with {dataset['Feedback'].min()}, and",
                    "size": 14,
                    "bold": False,
                    "ital": False},
                7: {"level": 0,
                    "align": "left",
                    "text": f"{tuple(dataset.loc[dataset[['Feedback']].idxmax(), 'City'])[0]} in {tuple(dataset.loc[dataset[['Feedback']].idxmax(), 'Country'])[0]} has the highest rating with {dataset['Feedback'].max()}.",
                    "size": 14,
                    "bold": False,
                    "ital": False}
                }
            }
        },
        4: {"layout": 2, "title": 'Some insights 2/5',
            "placeholders":{1:
                {1: {"level": 0,
                    "align": "left",
                    "text": "Here is the list of all cities having over 200 all-inclusive accommodations:",
                    "size": 14,
                    "bold": False,
                    "ital": False},
                2: {"level": 0,
                    "align": "left",
                    "text": f"{dataset.loc[(dataset['All-inc'] > 200), ['Country', 'City', 'All-inc']]}",
                    "size": 10,
                    "bold": True,
                    "ital": False}
                }
            }
        },
        5: {"layout": 2, "title": 'Some insights 3/5',
            "placeholders":{1:
                {1: {"level": 0,
                    "align": "left",
                    "text": "Here is the list of all cities having a feedback score superior to 8:",
                    "size": 14,
                    "bold": False,
                    "ital": False},
                2: {"level": 0,
                    "align": "left",
                    "text": f"{dataset.loc[(dataset['Feedback'] > 8), ['Country', 'City', 'Feedback']]}",
                    "size": 10,
                    "bold": True,
                    "ital": False}
                }
            }

        },
        6: {"layout": 2, "title": 'Some insights 4/5',
            "placeholders":{1:
                {1: {"level": 0,
                    "align": "left",
                    "text": "Here is the list of all cities having a feedback score inferior to 2:",
                    "size": 14,
                    "bold": False,
                    "ital": False},
                2: {"level": 0,
                    "align": "left",
                    "text": f"{dataset.loc[(dataset['Feedback'] < 2), ['Country', 'City', 'Feedback']]}",
                    "size": 10,
                    "bold": True,
                    "ital": False}
                }
            }
        },
        7: {"layout": 2, "title": 'Some insights 5/5',
            "placeholders":{1:
                {1: {"level": 0,
                    "align": "left",
                    "text": "The results displayed on the last 2 slides can be surprising, but the feedback index was composed taking into account various components:",
                    "size": 14,
                    "bold": False,
                    "ital": False},
                2: {"level": 1,
                    "align": "left",
                    "text": "The ratio of rental accommodations compared to the number of yearly tourists",
                    "size": 12,
                    "bold": False,
                    "ital": False},
                3: {"level": 2,
                    "align": "left",
                    "text": "The lower it is, the higher tourist interest",
                    "size": 10,
                    "bold": False,
                    "ital": True},
                4: {"level": 1,
                    "align": "left",
                    "text": "The ratio of rental accommodations compared to the number of inhabitants",
                    "size": 12,
                    "bold": False,
                    "ital": False},
                5: {"level": 2,
                    "align": "left",
                    "text": "The higher, the higher tourism as a share of the GDP",
                    "size": 10,
                    "bold": False,
                    "ital": True},
                6: {"level": 1,
                    "align": "left",
                    "text": "The ratio of all-inclusive accommodations compared to all rental accommodation",
                    "size": 12,
                    "bold": False,
                    "ital": False},
                7: {"level": 2,
                    "align": "left",
                    "text": "The higher, the higher client satisfaction",
                    "size": 10,
                    "bold": False,
                    "ital": True},
                8: {"level": 1,
                    "align": "left",
                    "text": "The ratio of the average number of stars per accommodations",
                    "size": 12,
                    "bold": False,
                    "ital": False},
                9: {"level": 2,
                    "align": "left",
                    "text": "The higher, the higher the quality and standard of accommodation and therefore of client satisfaction",
                    "size": 10,
                    "bold": False,
                    "ital": True},
                10: {"level": 1,
                    "align": "left",
                    "text": "The ratio based on the risk level for the city",
                    "size": 12,
                    "bold": False,
                    "ital": False},
                11: {"level": 2,
                    "align": "left",
                    "text": "The higher, the safer the city",
                    "size": 10,
                    "bold": False,
                    "ital": True},
                12: {"level": 1,
                    "align": "left",
                    "text": "A differential ratio between the risk rating of the country and of the city",
                    "size": 12,
                    "bold": False,
                    "ital": False},
                13: {"level": 2,
                    "align": "left",
                    "text": "A positive differential indicates a safer environment",
                    "size": 10,
                    "bold": False,
                    "ital": True},
                14: {"level": 0,
                    "align": "left",
                    "text": "This leads to some cities such as Merida in Venezuela having the best rating, and the Greek cities being amongst the lowest rating.",
                    "size": 14,
                    "bold": False,
                    "ital": False}
                }
            }

        },
        8: {"layout": 8, "title": 'Correlation search 1/4',
            "placeholders":{1: 'Comp_feedback_allInc.jpg',
                2: 'It seems that cities with a high level of all-inclusive accommodations (above 400) do tend to have a higher feedback rating (above 3), but this is not really significant as many cities with a lower number of accommodations end up with a high feedback mark and as the highest feedback mark belong to two cities with less than 200 all-inclusive accommodations.'}
        },
        9: {"layout": 8, "title": 'Correlation search 2/4',
            "placeholders":{1: 'Comp_feedback_stars.jpg',
                2: 'No correlation was found between the average star rating and the number of all-inclusive accommodations.'}
        },
        10: {"layout": 8, "title": 'Correlation search 3/4',
            "placeholders":{1: 'Bar_chart.jpg',
                2: 'No direct correlation was found the average star rating and the feedback score.'}
        },
        11: {"layout": 8, "title": 'Correlation search 4/4',
            "placeholders":{1: 'Comp_allInc_stars.jpg',
                2: 'No correlation was found between average number of stars and the number of all-inclusive accommodations.'}
        }

    }

    # use dictionary information to create the first 7 slides.
    for k, v in slide_content.items():
        sld(v, dataset)

    # saves the presentation
    pr1.save("TTA_HLA_Week_04_pandas.pptx")
    # and opens it
    os.startfile("TTA_HLA_Week_04_pandas.pptx")
    pass

if __name__ == '__main__':
    main()
